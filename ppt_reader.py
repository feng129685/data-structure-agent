"""
数据结构PPT知识提炼工具
提取PPT文本 → DeepSeek提炼 → 结构化知识库（适配AI智能体训练）
- 分层过滤无效内容（课堂口语、过渡话术等）
- 6类实体提取（定义、逻辑、代码、复杂度、易错点、例题）
- 知识图谱 + 逐页拆解 + 对比库 + 图文转译 + 问答对 + 错误诊断
"""
import streamlit as st
import os
from pptx import Presentation
from openai import OpenAI

st.set_page_config(page_title="PPT知识提炼", page_icon="🧠", layout="wide")

# ===== 配置 =====
PPT_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "data_struct_ppt_analys")
API_KEY = "sk-2e5b5b541fd842ef813841f753a6ef62"
BASE_URL = "https://api.deepseek.com"

CHAPTER_MAP = {
    '第一章': '第一章 绪论', '第二章': '第二章 线性表',
    '第三章': '第三章 栈和队列', '第四章': '第四章 串',
    '第五章': '第五章 数组和广义表', '第六章': '第六章 树和二叉树',
    '第七章': '第七章 图', '第八章': '第八章 查找', '第九章': '第九章 排序',
}

# ===== PPT内容提取 =====

def extract_ppt_content(ppt_path: str) -> dict:
    """提取PPT全部文本"""
    prs = Presentation(ppt_path)
    total = len(prs.slides)
    all_text_parts = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    slide_texts.append(t)

        if slide_texts:
            all_text_parts.append(f"【第{slide_num}页】\n" + "\n".join(slide_texts))

    return {
        'text': "\n\n---\n\n".join(all_text_parts),
        'total_slides': total,
    }


# ===== DeepSeek 提炼 =====

REFINE_PROMPT = """你是一位数据结构课程的知识提炼专家，专门为AI智能体训练和知识库构建提炼PPT内容。
请仔细阅读以下PPT课件内容（共{total}页），严格按照下述规则输出结构化提炼结果。

PPT内容：
---
{content}
---

## 提炼规则（必须严格遵守）

### 一、分层过滤
1. **完全丢弃**：课堂引导句、开场白、趣味科普、分页过渡文字、动画说明、配图修饰文字、互动提问话术、无关拓展小故事。例："下面我们一起来学习链表""大家看这张图是不是很直观"全部删除。
2. **压缩为备注**：历史发展、简单应用场景举例、硬件无关类比（如排队类比栈），仅用1行标注"📎拓展：..."。
3. **核心提取**：所有用于算法推理、代码编写、习题解答、考试答疑的内容，分为6类入库：
   - ①概念定义：结构本质、存储方式、适用场景
   - ②逻辑结构：节点组成、头尾指针、父子关系、链式/顺序存储区别
   - ③基础操作：增删改查、初始化、销毁、遍历、排序
   - ④算法复杂度：T(n)、S(n)、最好/最坏/平均
   - ⑤代码实现：伪代码、C/C++完整代码、边界处理
   - ⑥考点易错：空指针、内存泄漏、越界、循环条件错误、栈溢出

### 二、输出格式（严格按以下结构）

## 🏗️ 一、知识图谱（层级树）
（先提取PPT大标题，搭建层级树，如：
数据结构
├─线性结构
│  ├─顺序表
│  ├─单链表/双向链表/循环链表
│  ├─栈、队列
├─树形结构
│  ├─二叉树、二叉搜索树
│  ├─平衡树、堆
...
仅包含本PPT涉及的知识点分支，标注所在页码。）

## 📋 二、逐页标准化拆解
（每页按固定9字段提取，无则写"无"）

> 模板：
> 【页码】第X页
> 【所属模块】如「单链表 - 插入操作」
> 【核心定义】极简专业描述，1-2句话
> 【存储模型】顺序存储 / 链式存储 / 索引存储 / 散列存储
> 【算法流程】分步骤关键词，用 → 串联逻辑
> 【伪代码/代码】保留边界判断、循环终止条件
> 【复杂度】T(n)=... S(n)=... 最好/最坏/平均
> 【边界用例】空表、头插、尾插、中间插入等
> 【高频易错】标签标记，用于智能体纠错

## 🔄 三、对比知识库
（同类结构横向对比，如：
- 顺序表 VS 链表：增删 | 查询 | 内存开销 对比
- 栈 VS 队列：出入规则 | 应用场景 对比
）

## 📝 四、图文转译
（PPT中若有链表/树/栈/图等示意图，必须文字化描述其流程：
- 顺序表扩容：重新开辟2倍空间 → 复制旧数组 → 释放原内存
- 二叉树遍历：递归执行步骤、栈出入栈顺序
- 图邻接矩阵/邻接表：存储规则、查询邻点逻辑
无示意图则写"本PPT无示意图"。）

## ❓ 五、习题问答对
（所有例题、练习题拆为标准问答对：
- **Q**：题目描述
- **A**：算法步骤 + 完整代码 + 复杂度 + 边界分析
无习题则写"本PPT无习题"。）

## ⚠️ 六、错误诊断规则
（从PPT标红易错点、常见bug、考试陷阱提取，每条格式：
- **错误现象**：...
- **根因**：...
- **修复**：...
无则写"本PPT未涉及"。）

---

**输出原则**：
- 口语、类比全部替换为严谨计算机术语
- 复杂度必须独立成行，不混入段落
- 代码必须保留边界条件（空结构、长度为1、下标0）
- 知识点标注上下级关联，支持多轮推理跳转
- 只输出有实际内容的部分，无内容则写"无"即可，不编造"""


def analyze_with_deepseek(text: str, total_slides: int):
    """PPT内容喂给DeepSeek做知识提炼（流式）"""
    # 截断过长文本（保留前25000字符）
    content = text[:25000]
    if len(text) > 25000:
        content += f"\n\n...（原文共{len(text)}字符，已截取前25000字符）"

    prompt = REFINE_PROMPT.format(total=total_slides, content=content)

    client = OpenAI(api_key=API_KEY, base_url=BASE_URL)
    return client.chat.completions.create(
        model="deepseek-chat",
        messages=[
            {"role": "system", "content": "你是数据结构知识提炼专家，为AI智能体训练构建结构化知识库。严格按规则过滤、提取、结构化输出，不编造内容。"},
            {"role": "user", "content": prompt}
        ],
        stream=True,
        temperature=0.2,
        max_tokens=8000,
    )


# ===== Session State =====
def init():
    for k, v in {
        'selected_file': None, 'ppt_path': None,
        'analysis_done': False, 'analysis_result': '',
        'ppt_data': None,
    }.items():
        if k not in st.session_state:
            st.session_state[k] = v


# ===== 文件扫描 =====
def scan_files():
    if not os.path.exists(PPT_DIR):
        return []
    files = []
    for f in os.listdir(PPT_DIR):
        if f.endswith('.pptx'):
            ch = '其他'
            for k in CHAPTER_MAP:
                if k in f: ch = k; break
            files.append({
                'name': f, 'chapter': ch,
                'chapter_name': CHAPTER_MAP.get(ch, '其他'),
                'path': os.path.join(PPT_DIR, f),
                'size_mb': os.path.getsize(os.path.join(PPT_DIR, f)) / (1024 * 1024),
            })
    files.sort(key=lambda x: (list(CHAPTER_MAP).index(x['chapter']) if x['chapter'] in CHAPTER_MAP else 99, x['name']))
    return files


# ===== UI =====
def main():
    init()
    files = scan_files()

    if not files:
        st.error("未找到PPT文件")
        return

    # ---- 侧边栏 ----
    st.sidebar.title("📚 数据结构PPT")
    from itertools import groupby
    for ch, group in groupby(files, key=lambda f: f['chapter']):
        gl = list(group)
        st.sidebar.markdown(f"**{gl[0]['chapter_name']}** ({len(gl)}个)")
        for f in gl:
            label = f"📄 {f['name'].replace('.pptx','')[:30]}"
            if st.sidebar.button(label, key=f['name'], use_container_width=True):
                st.session_state.selected_file = f['name']
                st.session_state.ppt_path = f['path']
                st.session_state.analysis_done = False
                st.session_state.analysis_result = ''
                st.session_state.ppt_data = None
                st.rerun()

    st.sidebar.markdown("---")
    st.sidebar.caption(f"共 {len(files)} 个PPT | DeepSeek 知识提炼")

    # ---- 欢迎页 ----
    if not st.session_state.selected_file:
        st.title("📖 数据结构PPT·知识提炼")
        st.markdown("### PPT文本提取 → AI知识提炼 → 结构化知识库（适配智能体训练）")
        st.markdown("---")
        st.info("👈 从左侧选择一个PPT文件，点击「开始提炼」即可")
        return

    # ---- 已选文件 ----
    file_name = st.session_state.selected_file.replace('.pptx', '')
    st.title(f"📖 {file_name}")

    # 提取内容
    if st.session_state.ppt_data is None:
        with st.spinner("📤 提取PPT内容中..."):
            st.session_state.ppt_data = extract_ppt_content(st.session_state.ppt_path)
        st.rerun()

    data = st.session_state.ppt_data
    st.markdown(f"📊 **{data['total_slides']}** 页 | "
                f"**{len(data['text'])}** 字符")

    # ---- 提炼按钮 ----
    if not st.session_state.analysis_done:
        st.markdown("---")
        if st.button("🚀 开始知识提炼", type="primary", use_container_width=True):
            with st.spinner("🤖 DeepSeek 正在提炼知识..."):
                try:
                    stream = analyze_with_deepseek(data['text'], data['total_slides'])
                    placeholder = st.empty()
                    result = ""
                    for chunk in stream:
                        if chunk.choices[0].delta.content:
                            result += chunk.choices[0].delta.content
                            placeholder.markdown(result + "▌")
                    placeholder.markdown(result)
                    st.session_state.analysis_result = result
                    st.session_state.analysis_done = True
                    st.rerun()
                except Exception as e:
                    st.error(f"提炼失败: {e}")

    # ---- 展示结果 ----
    if st.session_state.analysis_done:
        st.markdown("---")

        tab1, tab2, tab3 = st.tabs([
            "🧠 知识提炼结果",
            "📚 轻量化知识库",
            "💻 代码生成素材",
        ])

        with tab1:
            st.markdown(st.session_state.analysis_result)
            st.download_button(
                "📥 下载提炼结果 (.md)",
                st.session_state.analysis_result,
                file_name=f"{file_name}_知识提炼.md",
                mime="text/markdown",
            )

        with tab2:
            st.markdown("### 📚 轻量化知识库（检索/答疑智能体适用）")
            st.markdown("从上方提炼结果中提取以下内容，按标签分类存储，支持关键词检索与知识点关联查询：")
            st.markdown("""
- **概念定义** → 标签：`定义` — 结构本质、存储方式、适用场景
- **逻辑结构** → 标签：`结构` — 节点组成、指针关系、存储区别
- **基础操作** → 标签：`操作` — 增删改查、初始化、遍历
- **算法复杂度** → 标签：`复杂度` — T(n)、S(n)、最好/最坏/平均
- **考点易错** → 标签：`易错` — 空指针、内存泄漏、越界等
            """)
            st.info("💡 将提炼结果按标签存储到向量数据库，即可构建可检索的知识库供答疑智能体调用。")

        with tab3:
            st.markdown("### 💻 代码生成素材库（代码编写智能体适用）")
            st.markdown("从上方提炼结果中提取以下内容，过滤理论段落，保留可执行代码信息：")
            st.markdown("""
- **操作流程**：步骤化描述，每步对应代码块
- **完整代码**：C/C++ 实现，含注释
- **边界处理**：空结构、长度为1、下标0等
- **复杂度标注**：每段代码附带 T(n)/S(n)
            """)
            st.info("💡 将操作流程+代码+边界处理配对存储，代码智能体可根据需求检索并输出可运行代码。")

        # 重新提炼
        if st.button("🔄 重新提炼"):
            st.session_state.analysis_done = False
            st.session_state.analysis_result = ''
            st.rerun()


if __name__ == "__main__":
    main()
